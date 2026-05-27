from PIL import Image
import pptx
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.util import Pt
from pptx.oxml.ns import qn
import io
import base64
import uuid
import hashlib
import re

# ── Helpers ─────────────────────────────────────────────────────────────────

_THEME_BUNDLE_CACHE = {}

def _uid(prefix="el"):
    return f"pptx_{prefix}_{uuid.uuid4().hex[:8]}"


def _compress_image(image_blob: bytes):
    """Compress to JPEG ≤ 1280×720 for reasonable SCORM package sizes."""
    try:
        with Image.open(io.BytesIO(image_blob)) as img:
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            img.thumbnail((1280, 720), Image.Resampling.LANCZOS)
            out = io.BytesIO()
            img.save(out, format="JPEG", quality=75)
            return out.getvalue(), "image/jpeg"
    except Exception:
        return image_blob, "image/png"


def _blob_to_data_uri(blob: bytes, mime: str) -> str:
    return f"data:{mime};base64,{base64.b64encode(blob).decode()}"


def _rgb_to_hex(rgb) -> str:
    """Convert pptx RGBColor (or tuple) to CSS hex string."""
    try:
        return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
    except Exception:
        pass
    try:
        r, g, b = rgb
        return f"#{r:02x}{g:02x}{b:02x}"
    except Exception:
        return ""


def _theme_cache_key(theme_part) -> str:
    try:
        return str(theme_part.partname)
    except Exception:
        return str(id(theme_part))


def _get_theme_part(part):
    if part is None:
        return None
    for rel in part.rels.values():
        if "relationships/theme" in rel.reltype:
            return rel.target_part
    return None


def _extract_theme_bundle(theme_part) -> dict:
    """
    Extract theme colors plus theme background fill styles.

    `p:bgRef` values point into the theme's background fill style list, so we
    cache both the scheme colors and the raw bg fill definitions.
    """
    bundle = {"colors": {}, "bg_fill_styles": {}}
    if theme_part is None:
        return bundle

    try:
        if hasattr(theme_part, "element"):
            root = theme_part.element
        else:
            from pptx.oxml import parse_xml
            root = parse_xml(theme_part.blob)

        theme_elements = root.find(qn("a:themeElements"))
        if theme_elements is None:
            theme_elements = root.find(f".//{qn('a:themeElements')}")
        if theme_elements is None:
            return bundle

        clr_scheme = theme_elements.find(qn("a:clrScheme"))
        if clr_scheme is None:
            clr_scheme = theme_elements.find(f".//{qn('a:clrScheme')}")

        if clr_scheme is not None:
            for child in clr_scheme:
                tag_name = child.tag.split("}")[-1]

                srgb_clr = child.find(qn("a:srgbClr"))
                if srgb_clr is not None:
                    val = (srgb_clr.get("val") or "").strip()
                    if re.fullmatch(r"[0-9A-Fa-f]{6}", val):
                        bundle["colors"][tag_name] = f"#{val.lower()}"
                    continue

                sys_clr = child.find(qn("a:sysClr"))
                if sys_clr is not None:
                    val = (sys_clr.get("lastClr") or "").strip()
                    if re.fullmatch(r"[0-9A-Fa-f]{6}", val):
                        bundle["colors"][tag_name] = f"#{val.lower()}"
                    else:
                        sys_val = (sys_clr.get("val") or "").strip().lower()
                        if sys_val == "windowtext":
                            bundle["colors"][tag_name] = "#000000"
                        elif sys_val == "window":
                            bundle["colors"][tag_name] = "#ffffff"

        fmt_scheme = theme_elements.find(qn("a:fmtScheme"))
        if fmt_scheme is None:
            fmt_scheme = theme_elements.find(f".//{qn('a:fmtScheme')}")

        if fmt_scheme is not None:
            bg_fill_style_lst = fmt_scheme.find(qn("a:bgFillStyleLst"))
            if bg_fill_style_lst is not None:
                for idx, fill_el in enumerate(list(bg_fill_style_lst), start=1001):
                    bundle["bg_fill_styles"][idx] = fill_el

    except Exception as e:
        print(f"[pptx_parser] Could not extract theme bundle: {e}")

    return bundle


def _get_cached_theme_bundle(theme_part) -> dict:
    if theme_part is None:
        return {"colors": {}, "bg_fill_styles": {}}
    cache_key = _theme_cache_key(theme_part)
    if cache_key not in _THEME_BUNDLE_CACHE:
        _THEME_BUNDLE_CACHE[cache_key] = _extract_theme_bundle(theme_part)
    return _THEME_BUNDLE_CACHE[cache_key]


def _get_theme_part_for_slide(slide):
    try:
        master = getattr(getattr(slide, "slide_layout", None), "slide_master", None)
        if master is not None:
            return _get_theme_part(master.part)
    except Exception:
        pass
    return None


def _build_theme_color_map(prs) -> dict:
    """
    Extract actual hex colors from the presentation's theme XML.
    Tries every slide master so multi-master presentations are covered.
    Returns a map of scheme name → CSS hex, e.g. {"accent1": "#c0392b", ...}
    """
    theme_colors = {}
    try:
        masters = list(prs.slide_masters)
        if not masters:
            return theme_colors

        for master in masters:
            try:
                theme_part = _get_theme_part(master.part)
            except Exception:
                continue

            bundle = _get_cached_theme_bundle(theme_part)
            for tag_name, hex_color in bundle.get("colors", {}).items():
                if tag_name not in theme_colors:
                    theme_colors[tag_name] = hex_color

    except Exception as e:
        print(f"[pptx_parser] Could not extract theme colors: {e}")

    return theme_colors


def _extract_color_map_attrs(color_map_el) -> dict:
    """Return a PresentationML clrMap/overrideClrMapping element as a dict."""
    if color_map_el is None:
        return {}
    keys = (
        "bg1", "tx1", "bg2", "tx2",
        "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
        "hlink", "folHlink",
    )
    return {
        key: value
        for key in keys
        if (value := (color_map_el.get(key) or "").strip())
    }


def _get_effective_color_map(slide) -> dict:
    """
    Resolve the effective scheme-color mapping for a slide.

    OOXML scheme colors like `tx1` and `bg1` are indirections; the actual
    theme slot comes from the slide/master clrMap chain.
    """
    try:
        color_map = {}

        master_el = _xml_element(slide.slide_layout.slide_master)
        if master_el is not None:
            color_map.update(_extract_color_map_attrs(master_el.find(qn("p:clrMap"))))

        for owner in (getattr(slide, "slide_layout", None), slide):
            if owner is None:
                continue
            owner_el = _xml_element(owner)
            if owner_el is None:
                continue
            clr_map_ovr = owner_el.find(qn("p:clrMapOvr"))
            if clr_map_ovr is None:
                continue
            override = clr_map_ovr.find(qn("p:overrideClrMapping")) or clr_map_ovr.find(qn("a:overrideClrMapping"))
            if override is not None:
                color_map.update(_extract_color_map_attrs(override))
            elif (
                clr_map_ovr.find(qn("p:masterClrMapping")) is not None
                or clr_map_ovr.find(qn("a:masterClrMapping")) is not None
            ):
                # Explicitly keep the master mapping as-is.
                pass

        return color_map
    except Exception:
        return {}


def _resolve_scheme_key(scheme_key: str, color_map=None) -> str:
    """Map a scheme key like `tx1` through the effective clrMap if needed."""
    if not scheme_key:
        return ""
    resolved = scheme_key
    if color_map:
        resolved = color_map.get(scheme_key, scheme_key)

    # Fallback to standard OOXML mapping if not overridden or maps to self
    if resolved in ("bg1", "tx1", "bg2", "tx2"):
        return {
            "bg1": "lt1",
            "tx1": "dk1",
            "bg2": "lt2",
            "tx2": "dk2",
        }[resolved]

    return resolved


def _xml_element(obj):
    """Return the underlying lxml element for a python-pptx proxy object."""
    return getattr(obj, "element", None) or getattr(obj, "_element", None)


def _apply_color_transforms(hex_color: str, scheme_el) -> str:
    """
    Apply lum/shade/tint modifiers from a <a:schemeClr> element to the
    resolved hex color.  Returns the transformed hex string.
    """
    if not hex_color or len(hex_color) != 7:
        return hex_color
    try:
        r = int(hex_color[1:3], 16) / 255.0
        g = int(hex_color[3:5], 16) / 255.0
        b = int(hex_color[5:7], 16) / 255.0

        for child in scheme_el:
            local = child.tag.split('}')[-1]
            val_str = child.get("val", "0")
            try:
                val = int(val_str) / 100000.0  # OOXML uses 1/1000ths of a percent
            except ValueError:
                continue

            if local == "lumMod":
                r, g, b = r * val, g * val, b * val
            elif local == "lumOff":
                r, g, b = r + val, g + val, b + val
            elif local == "shade":
                r, g, b = r * val, g * val, b * val
            elif local == "tint":
                r = r + (1.0 - r) * val
                g = g + (1.0 - g) * val
                b = b + (1.0 - b) * val

        r = max(0, min(1, r))
        g = max(0, min(1, g))
        b = max(0, min(1, b))
        return f"#{int(r*255):02x}{int(g*255):02x}{int(b*255):02x}"
    except Exception:
        return hex_color


def _extract_hex_from_color_element(color_el, theme_map=None, color_map=None, placeholder_hex="") -> str:
    """Resolve a single OOXML color node such as schemeClr or srgbClr."""
    if color_el is None:
        return ""

    local = color_el.tag.split("}")[-1] if "}" in color_el.tag else color_el.tag

    if local == "srgbClr":
        val = (color_el.get("val") or "").strip()
        if re.fullmatch(r"[0-9A-Fa-f]{6}", val):
            return _apply_color_transforms(f"#{val.lower()}", color_el)
        return ""

    if local == "scrgbClr":
        try:
            r = max(0, min(255, round(int(color_el.get("r", "0")) * 255 / 100000)))
            g = max(0, min(255, round(int(color_el.get("g", "0")) * 255 / 100000)))
            b = max(0, min(255, round(int(color_el.get("b", "0")) * 255 / 100000)))
            return _apply_color_transforms(f"#{r:02x}{g:02x}{b:02x}", color_el)
        except Exception:
            return ""

    if local == "schemeClr":
        val = (color_el.get("val") or "").strip()
        if val == "phClr" and placeholder_hex:
            resolved = placeholder_hex
        else:
            val = _resolve_scheme_key(val, color_map)
            resolved = theme_map.get(val, "") if theme_map else ""
        if resolved:
            resolved = _apply_color_transforms(resolved, color_el)
        return resolved

    if local == "sysClr":
        val = (color_el.get("lastClr") or "").strip()
        if re.fullmatch(r"[0-9A-Fa-f]{6}", val):
            return f"#{val.lower()}"
        sys_val = (color_el.get("val") or "").strip().lower()
        if sys_val == "windowtext":
            return "#000000"
        elif sys_val == "window":
            return "#ffffff"
        return ""

    if local == "prstClr":
        preset_colors = {
            "black": "#000000", "white": "#ffffff", "red": "#ff0000",
            "green": "#008000", "blue": "#0000ff", "yellow": "#ffff00",
            "cyan": "#00ffff", "magenta": "#ff00ff", "orange": "#ffa500",
            "purple": "#800080", "brown": "#a52a2a", "gray": "#808080",
            "grey": "#808080", "pink": "#ffc0cb", "navy": "#000080",
            "teal": "#008080", "lime": "#00ff00", "maroon": "#800000",
            "olive": "#808000", "silver": "#c0c0c0", "aqua": "#00ffff",
            "fuchsia": "#ff00ff",
        }
        val = (color_el.get("val") or "").strip().lower()
        return preset_colors.get(val, "")

    return ""


def _extract_hex_from_color_choice(color_parent, theme_map=None, color_map=None, placeholder_hex="") -> str:
    """
    Extract a CSS hex color from a pptx color XML element.

    Resolution order:
      1. <a:srgbClr> — explicit RGB, used directly.
      2. <a:schemeClr> — look up in the *dynamic* theme_map extracted from
         the actual presentation theme.  The hardcoded fallback map has been
         removed because it caused brown→green and red→green color corruption
         when a presentation's accent colors differed from the old defaults.
         If the scheme key is not found in theme_map we return "" so the
         caller can fall through gracefully (background → "none").
      3. <a:sysClr> — system color, use lastClr attribute.
      4. <a:prstClr> — preset color name, converted to hex.
    """
    if color_parent is None:
        return ""

    for tag in ("srgbClr", "scrgbClr", "schemeClr", "sysClr", "prstClr"):
        try:
            color_el = color_parent.find(qn(f"a:{tag}"))
            if color_el is not None:
                return _extract_hex_from_color_element(
                    color_el,
                    theme_map=theme_map,
                    color_map=color_map,
                    placeholder_hex=placeholder_hex,
                )
        except Exception:
            pass

    return ""


def _emu_to_pct(val, total) -> float:
    """Convert EMU position/size to a percentage of the slide dimension."""
    try:
        return round(val / total * 100, 2)
    except Exception:
        return 0.0


# ── Rich text extraction ─────────────────────────────────────────────────────

def _extract_color_from_rpr(r_pr, theme_map=None, color_map=None) -> str:
    """Resolve a color from a run-properties-like element when present."""
    if r_pr is None:
        return ""
    try:
        solid_fill = r_pr.find(qn("a:solidFill"))
        if solid_fill is not None:
            return _extract_hex_from_color_choice(solid_fill, theme_map, color_map)
    except Exception:
        pass
    return ""


def _placeholder_style_tag(shape) -> str:
    """Best-effort txStyles bucket for a shape."""
    try:
        if getattr(shape, "is_placeholder", False):
            ph_type = getattr(shape.placeholder_format.type, "name", "")
            if ph_type in {"TITLE", "CENTER_TITLE"}:
                return "titleStyle"
            if ph_type in {"BODY", "SUBTITLE", "OBJECT"}:
                return "bodyStyle"
    except Exception:
        pass
    return "bodyStyle"


def _extract_inherited_text_hex_color(run, para, shape, slide, theme_map=None, color_map=None) -> str:
    """
    Resolve inherited text color from paragraph/list/master text styles.

    PowerPoint often stores visible text color above the run itself, especially
    for placeholders and master-driven themes.
    """
    try:
        p_pr = para._p.find(qn("a:pPr"))
        if p_pr is not None:
            hex_color = _extract_color_from_rpr(p_pr.find(qn("a:defRPr")), theme_map, color_map)
            if hex_color:
                return hex_color
    except Exception:
        pass

    try:
        tx_body = shape.text_frame._txBody
        lst_style = tx_body.find(qn("a:lstStyle"))
        if lst_style is not None:
            lvl = 0
            try:
                lvl = para.level
            except Exception:
                pass
            lvl_tag = qn(f"a:lvl{max(1, min(9, lvl + 1))}pPr")
            lvl_ppr = lst_style.find(lvl_tag)
            if lvl_ppr is not None:
                hex_color = _extract_color_from_rpr(lvl_ppr.find(qn("a:defRPr")), theme_map, color_map)
                if hex_color:
                    return hex_color
            def_ppr = lst_style.find(qn("a:defPPr"))
            if def_ppr is not None:
                hex_color = _extract_color_from_rpr(def_ppr.find(qn("a:defRPr")), theme_map, color_map)
                if hex_color:
                    return hex_color
    except Exception:
        pass

    try:
        if getattr(shape, "is_placeholder", False):
            placeholder_idx = getattr(shape.placeholder_format, "idx", None)
            for layout_shape in slide.slide_layout.shapes:
                if not getattr(layout_shape, "is_placeholder", False):
                    continue
                if getattr(layout_shape.placeholder_format, "idx", None) != placeholder_idx:
                    continue
                if not hasattr(layout_shape, "text_frame"):
                    continue

                tx_body = layout_shape.text_frame._txBody
                lst_style = tx_body.find(qn("a:lstStyle"))
                if lst_style is not None:
                    lvl = 0
                    try:
                        lvl = para.level
                    except Exception:
                        pass
                    lvl_tag = qn(f"a:lvl{max(1, min(9, lvl + 1))}pPr")
                    lvl_ppr = lst_style.find(lvl_tag)
                    if lvl_ppr is not None:
                        hex_color = _extract_color_from_rpr(lvl_ppr.find(qn("a:defRPr")), theme_map, color_map)
                        if hex_color:
                            return hex_color
                    def_ppr = lst_style.find(qn("a:defPPr"))
                    if def_ppr is not None:
                        hex_color = _extract_color_from_rpr(def_ppr.find(qn("a:defRPr")), theme_map, color_map)
                        if hex_color:
                            return hex_color
                break
    except Exception:
        pass

    try:
        master_el = _xml_element(slide.slide_layout.slide_master)
        tx_styles = master_el.find(qn("p:txStyles")) if master_el is not None else None
        if tx_styles is not None:
            style_tag = qn(f"p:{_placeholder_style_tag(shape)}")
            style_el = tx_styles.find(style_tag)
            if style_el is not None:
                lvl = 0
                try:
                    lvl = para.level
                except Exception:
                    pass
                lvl_tag = qn(f"a:lvl{max(1, min(9, lvl + 1))}pPr")
                lvl_ppr = style_el.find(lvl_tag)
                if lvl_ppr is not None:
                    hex_color = _extract_color_from_rpr(lvl_ppr.find(qn("a:defRPr")), theme_map, color_map)
                    if hex_color:
                        return hex_color
                def_ppr = style_el.find(qn("a:defPPr"))
                if def_ppr is not None:
                    hex_color = _extract_color_from_rpr(def_ppr.find(qn("a:defRPr")), theme_map, color_map)
                    if hex_color:
                        return hex_color
    except Exception:
        pass

    return ""


def _extract_run_hex_color(run, para, shape, slide, theme_map=None, color_map=None) -> str:
    """
    Resolve a run's font color from OOXML first so theme colors use the
    presentation's actual palette rather than python-pptx's lossy `.rgb` path.
    """
    try:
        r_pr = run._r.find(qn("a:rPr"))
        hex_color = _extract_color_from_rpr(r_pr, theme_map, color_map)
        if hex_color:
            return hex_color
    except Exception:
        pass

    try:
        color_format = run.font.color
        if color_format is not None and getattr(color_format, "type", None) == MSO_COLOR_TYPE.SCHEME:
            theme_color = getattr(color_format, "theme_color", None)
            theme_name = getattr(theme_color, "name", "")
            if theme_name and theme_map:
                theme_key = {
                    "TEXT_1": "tx1",
                    "TEXT_2": "tx2",
                    "BACKGROUND_1": "bg1",
                    "BACKGROUND_2": "bg2",
                    "DARK_1": "dk1",
                    "DARK_2": "dk2",
                    "LIGHT_1": "lt1",
                    "LIGHT_2": "lt2",
                    "FOLLOWED_HYPERLINK": "folHlink",
                    "HYPERLINK": "hlink",
                }.get(theme_name, theme_name.lower().replace("_", ""))
                theme_key = _resolve_scheme_key(theme_key, color_map)
                hex_color = theme_map.get(theme_key, "")
                if hex_color:
                    return hex_color
    except Exception:
        pass

    try:
        color = run.font.color.rgb
        if color:
            return _rgb_to_hex(color)
    except Exception:
        pass

    hex_color = _extract_inherited_text_hex_color(run, para, shape, slide, theme_map, color_map)
    if hex_color:
        return hex_color

    return ""


def _para_to_html(para, shape, slide, theme_map=None, color_map=None) -> str:
    """Convert a pptx paragraph to inline HTML, preserving basic formatting."""
    parts = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        style_parts = []
        try:
            if run.font.bold:
                style_parts.append("font-weight:bold")
        except Exception:
            pass
        try:
            if run.font.italic:
                style_parts.append("font-style:italic")
        except Exception:
            pass
        try:
            if run.font.underline:
                style_parts.append("text-decoration:underline")
        except Exception:
            pass
        try:
            size = run.font.size
            if size:
                style_parts.append(f"font-size:{int(size / 12700)}px")
        except Exception:
            pass
        hex_color = _extract_run_hex_color(run, para, shape, slide, theme_map, color_map)
        if hex_color:
            style_parts.append(f"color:{hex_color}")

        escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if style_parts:
            parts.append(f'<span style="{"; ".join(style_parts)}">{escaped}</span>')
        else:
            parts.append(escaped)

    return "".join(parts)


def _is_numbered_list(para) -> bool:
    """Detect if a paragraph is in a numbered list via pptx XML."""
    try:
        pPr = para._p.find(qn("a:pPr"))
        if pPr is not None:
            buAutoNum = pPr.find(qn("a:buAutoNum"))
            if buAutoNum is not None:
                return True
    except Exception:
        pass
    return False


def _is_bullet_list(para) -> bool:
    """Detect if a paragraph has a bullet character (non-auto-num)."""
    try:
        pPr = para._p.find(qn("a:pPr"))
        if pPr is not None:
            buNone = pPr.find(qn("a:buNone"))
            if buNone is not None:
                return False
            buChar = pPr.find(qn("a:buChar"))
            buFont = pPr.find(qn("a:buFont"))
            if buChar is not None or buFont is not None:
                return True
    except Exception:
        pass
    return False


def _get_indent_level(para) -> int:
    try:
        pPr = para._p.find(qn("a:pPr"))
        if pPr is not None:
            lvl = pPr.get("lvl")
            if lvl:
                return int(lvl)
    except Exception:
        pass
    return 0


# ── Background extraction ────────────────────────────────────────────────────

def _hex_to_rgb_tuple(hex_color: str) -> tuple[int, int, int] | None:
    if not hex_color or len(hex_color) != 7 or not hex_color.startswith("#"):
        return None
    try:
        return (
            int(hex_color[1:3], 16),
            int(hex_color[3:5], 16),
            int(hex_color[5:7], 16),
        )
    except Exception:
        return None


def _interpolate_hex_colors(start_hex: str, end_hex: str, ratio: float) -> str:
    start_rgb = _hex_to_rgb_tuple(start_hex)
    end_rgb = _hex_to_rgb_tuple(end_hex)
    if start_rgb is None:
        return end_hex
    if end_rgb is None:
        return start_hex
    ratio = max(0.0, min(1.0, ratio))
    r = round(start_rgb[0] + (end_rgb[0] - start_rgb[0]) * ratio)
    g = round(start_rgb[1] + (end_rgb[1] - start_rgb[1]) * ratio)
    b = round(start_rgb[2] + (end_rgb[2] - start_rgb[2]) * ratio)
    return f"#{r:02x}{g:02x}{b:02x}"


def _extract_gradient_midpoint_color(grad_fill, theme_map=None, color_map=None, placeholder_hex="") -> str:
    """
    Collapse a gradient into a representative midpoint color.

    CourseForge currently stores one slide background color, so midpoint
    sampling is the closest single-color approximation we can preserve.
    """
    if grad_fill is None:
        return ""

    try:
        gs_lst = grad_fill.find(qn("a:gsLst"))
        if gs_lst is None:
            return ""

        stops = []
        for gs in gs_lst.findall(qn("a:gs")):
            hex_color = _extract_hex_from_color_choice(
                gs,
                theme_map=theme_map,
                color_map=color_map,
                placeholder_hex=placeholder_hex,
            )
            if not hex_color:
                continue
            try:
                pos = int(gs.get("pos", "0"))
            except Exception:
                pos = 0
            stops.append((max(0, min(100000, pos)), hex_color))

        if not stops:
            return ""
        if len(stops) == 1:
            return stops[0][1]

        stops.sort(key=lambda item: item[0])
        target = 50000
        if target <= stops[0][0]:
            return stops[0][1]

        for (start_pos, start_hex), (end_pos, end_hex) in zip(stops, stops[1:]):
            if start_pos <= target <= end_pos:
                if end_pos == start_pos:
                    return end_hex
                ratio = (target - start_pos) / (end_pos - start_pos)
                return _interpolate_hex_colors(start_hex, end_hex, ratio)

        return stops[-1][1]
    except Exception:
        return ""


def _extract_image_fill(blip_fill, owner) -> dict:
    if blip_fill is None or owner is None:
        return {"type": "none", "value": None}
    try:
        blip = blip_fill.find(qn("a:blip"))
        if blip is None:
            return {"type": "none", "value": None}
        r_id = blip.get(qn("r:embed"))
        if not r_id:
            return {"type": "none", "value": None}
        part = owner.part.related_parts.get(r_id)
        if not part:
            return {"type": "none", "value": None}
        blob, mime = _compress_image(part.blob)
        return {"type": "image", "value": _blob_to_data_uri(blob, mime)}
    except Exception:
        return {"type": "none", "value": None}


def _extract_background_from_fill(fill_el, owner=None, theme_map=None, color_map=None, placeholder_hex="") -> dict:
    if fill_el is None:
        return {"type": "none", "value": None}

    local = fill_el.tag.split("}")[-1] if "}" in fill_el.tag else fill_el.tag

    if local == "solidFill":
        hex_color = _extract_hex_from_color_choice(
            fill_el,
            theme_map=theme_map,
            color_map=color_map,
            placeholder_hex=placeholder_hex,
        )
        if hex_color:
            return {"type": "color", "value": hex_color}

    elif local == "gradFill":
        hex_color = _extract_gradient_midpoint_color(
            fill_el,
            theme_map=theme_map,
            color_map=color_map,
            placeholder_hex=placeholder_hex,
        )
        if hex_color:
            return {"type": "color", "value": hex_color}

    elif local == "blipFill":
        return _extract_image_fill(fill_el, owner)

    return {"type": "none", "value": None}


def _extract_background_from_owner(owner, theme_map=None, color_map=None, theme_bundle=None) -> dict:
    owner_el = _xml_element(owner)
    if owner_el is None:
        return {"type": "none", "value": None}

    try:
        c_sld = owner_el.find(qn("p:cSld"))
        bg_el = (c_sld.find(qn("p:bg")) if c_sld is not None else None) or owner_el.find(qn("p:bg"))
        if bg_el is None:
            return {"type": "none", "value": None}

        bg_pr = bg_el.find(qn("p:bgPr"))
        if bg_pr is not None:
            for fill_tag in ("solidFill", "gradFill", "blipFill"):
                fill_el = bg_pr.find(qn(f"a:{fill_tag}"))
                resolved = _extract_background_from_fill(
                    fill_el,
                    owner=owner,
                    theme_map=theme_map,
                    color_map=color_map,
                )
                if resolved["type"] != "none" and resolved["value"]:
                    return resolved

        bg_ref = bg_el.find(qn("p:bgRef"))
        if bg_ref is not None:
            placeholder_hex = _extract_hex_from_color_choice(
                bg_ref,
                theme_map=theme_map,
                color_map=color_map,
            )
            try:
                idx = int(bg_ref.get("idx", ""))
            except Exception:
                idx = None

            if idx is not None and theme_bundle:
                fill_el = theme_bundle.get("bg_fill_styles", {}).get(idx)
                resolved = _extract_background_from_fill(
                    fill_el,
                    owner=owner,
                    theme_map=theme_map,
                    color_map=color_map,
                    placeholder_hex=placeholder_hex,
                )
                if resolved["type"] != "none" and resolved["value"]:
                    return resolved

            if placeholder_hex:
                return {"type": "color", "value": placeholder_hex}
    except Exception:
        pass

    return {"type": "none", "value": None}


def _extract_slide_background(slide, theme_map=None, color_map=None, theme_bundle=None) -> dict:
    """
    Return background info:  { "type": "color"|"image"|"none", "value": ... }
    """
    owners = [
        slide,
        getattr(slide, "slide_layout", None),
        getattr(getattr(slide, "slide_layout", None), "slide_master", None),
    ]

    for owner in owners:
        if owner is None:
            continue
        resolved = _extract_background_from_owner(
            owner,
            theme_map=theme_map,
            color_map=color_map,
            theme_bundle=theme_bundle,
        )
        if resolved["type"] != "none" and resolved["value"]:
            return resolved

    return {"type": "none", "value": None}


# ── Transition extraction ────────────────────────────────────────────────────

_TRANSITION_NAMES = {
    "blinds": "Blinds", "checker": "Checker", "circle": "Circle",
    "comb": "Comb", "cover": "Cover", "cut": "Cut", "diamond": "Diamond",
    "dissolve": "Dissolve", "fade": "Fade", "newsflash": "Newsflash",
    "plus": "Plus", "pull": "Pull", "push": "Push", "random": "Random",
    "randomBar": "Random Bar", "split": "Split", "strips": "Strips",
    "wedge": "Wedge", "wheel": "Wheel", "wipe": "Wipe", "zoom": "Zoom",
    "fly": "Fly Through", "glitter": "Glitter", "honeycomb": "Honeycomb",
    "morph": "Morph", "origami": "Origami", "pan": "Pan", "peel": "Peel",
    "prestige": "Prestige", "reveal": "Reveal", "ripple": "Ripple",
    "rotate": "Rotate", "shred": "Shred", "switch": "Switch",
    "vortex": "Vortex", "wind": "Wind",
}

def _extract_transition(slide) -> str | None:
    """Return human-readable transition name or None."""
    try:
        spTree = slide._element
        transition_el = spTree.find(f".//{qn('p:transition')}")
        if transition_el is None:
            transition_el = slide._element.find(qn("p:transition"))
        if transition_el is not None:
            for child in transition_el:
                local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                name = _TRANSITION_NAMES.get(local, local if local else None)
                if name:
                    return name
            return "Transition"
    except Exception:
        pass
    return None


# ── Audio extraction ─────────────────────────────────────────────────────────

_AUDIO_MIMES = {
    "mp3": "audio/mpeg", "wav": "audio/wav", "ogg": "audio/ogg",
    "m4a": "audio/mp4", "aac": "audio/aac", "flac": "audio/flac",
    "wma": "audio/x-ms-wma",
}

_AUDIO_XML_TAGS = {"audioFile", "wavAudioFile", "audioCd"}


def _find_audio_rid(element) -> str | None:
    for el in element.iter():
        local = el.tag.split("}")[-1] if "}" in el.tag else el.tag
        if local in _AUDIO_XML_TAGS:
            rid = el.get(qn("r:link")) or el.get(qn("r:embed"))
            if rid:
                return rid
    return None


def _try_extract_audio_from_shape(shape, slide, slide_w, slide_h,
                                   audio_blocks: list, seen_rids: set) -> None:
    try:
        rid = _find_audio_rid(shape._element)
        if not rid or rid in seen_rids:
            return
        seen_rids.add(rid)

        rel = slide.part.rels.get(rid)

        if rel is None or getattr(rel, "is_external", False):
            print(
                f"[pptx_parser] audio '{getattr(shape, 'name', '?')}' on slide "
                f"uses a linked (external) file (rId={rid}). "
                f"Only embedded audio is supported — re-save the PPTX with "
                f"'Embed media in file' enabled to include this track."
            )
            return

        part = rel.target_part

        ext = part.partname.rsplit(".", 1)[-1].lower()
        mime = _AUDIO_MIMES.get(ext, "audio/mpeg")
        blob = part.blob

        _MAX_INLINE_BYTES = 2 * 1024 * 1024  # 2 MB
        if len(blob) > _MAX_INLINE_BYTES:
            print(
                f"[pptx_parser] audio '{getattr(shape, 'name', '?')}' is "
                f"{len(blob) // 1024} KB. Embedding as a data URI will inflate "
                f"the SCORM JSON payload. Consider using file-based export "
                f"(pass scorm_output_dir to extract_slides) for audio files "
                f"larger than {_MAX_INLINE_BYTES // 1024} KB."
            )

        media_id = hashlib.md5(blob).hexdigest()[:12]

        position = {}
        try:
            position = {
                "x":      _emu_to_pct(shape.left,   slide_w),
                "y":      _emu_to_pct(shape.top,    slide_h),
                "width":  _emu_to_pct(shape.width,  slide_w),
                "height": _emu_to_pct(shape.height, slide_h),
            }
        except Exception as pos_err:
            print(f"[pptx_parser] could not read position for audio shape "
                  f"'{getattr(shape, 'name', '?')}': {pos_err}")

        audio_blocks.append({
            "id":        _uid("audio"),
            "type":      "audio",
            "label":     getattr(shape, "name", None) or "Audio Track",
            "audioUrl":  _blob_to_data_uri(blob, mime),
            "mediaId":   media_id,
            "autoPlay":  True,
            "loop":      False,
            "controls":  True,
            "mandatory":  False,
            "position":   position,
        })

    except Exception as err:
        print(f"[pptx_parser] failed to extract audio from shape "
              f"'{getattr(shape, 'name', '?')}': {err}")


def _extract_embedded_audio(slide, slide_w: int, slide_h: int) -> list:
    audio_blocks: list = []
    seen_rids: set = set()

    try:
        for shape in slide.shapes:
            _try_extract_audio_from_shape(
                shape, slide, slide_w, slide_h, audio_blocks, seen_rids
            )
    except Exception as err:
        print(f"[pptx_parser] slide-level audio scan failed: {err}")

    return audio_blocks


# ── Table extraction ─────────────────────────────────────────────────────────

def _extract_table(shape) -> dict | None:
    """Convert a pptx table shape to a CourseForge table block."""
    try:
        table = shape.table
        headers = []
        rows = []
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text_frame.text.strip() if cell.text_frame else "" for cell in row.cells]
            if row_idx == 0:
                headers = cells
            else:
                rows.append(cells)
        if not headers:
            return None
        return {
            "id": _uid("table"),
            "type": "table",
            "headers": headers,
            "rows": rows,
        }
    except Exception:
        return None


# ── Shape-level text parsing (rich text, lists) ──────────────────────────────

def _extract_text_shapes(shape, slide, theme_map=None, color_map=None) -> list:
    if not hasattr(shape, "text_frame"):
        return []

    tf = shape.text_frame
    paragraphs = [p for p in tf.paragraphs if p.text.strip()]
    if not paragraphs:
        return []

    is_title = (
        getattr(shape, "is_placeholder", False)
        and shape.placeholder_format is not None
        and shape.placeholder_format.idx == 0
    )
    if is_title:
        return []

    is_all_bullets = all(_is_bullet_list(p) or _is_numbered_list(p) for p in paragraphs)

    if is_all_bullets and len(paragraphs) > 1:
        items = [p.text.strip() for p in paragraphs if p.text.strip()]
        return [{
            "id": _uid("list"),
            "type": "list",
            "items": items,
        }]

    html_parts = []
    for para in paragraphs:
        para_html = _para_to_html(para, shape, slide, theme_map, color_map)
        if not para_html.strip():
            continue
        try:
            sz = para.runs[0].font.size if para.runs else None
            if sz and sz / 12700 >= 28:
                html_parts.append(f"<h2 style='margin:0 0 8px'>{para_html}</h2>")
                continue
        except Exception:
            pass
        html_parts.append(f"<p style='margin:0 0 6px'>{para_html}</p>")

    if not html_parts:
        return []

    return [{
        "id": _uid("text"),
        "type": "text",
        "content": "".join(html_parts),
    }]


# ── Speaker notes ────────────────────────────────────────────────────────────

def _extract_notes(slide) -> str | None:
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        text = tf.text.strip()
        if text:
            return text
    except Exception:
        pass
    return None


def _extract_shape_fill_color(shape, theme_map=None, color_map=None) -> str:
    """Extract standard shape fill color, mapping theme/scheme colors as needed."""
    try:
        if not hasattr(shape, "fill") or shape.fill is None:
            return ""
        fill = shape.fill
        fill_type = fill.type
        if fill_type is not None and str(fill_type) in ("SOLID", "1"):
            try:
                color = fill.fore_color.rgb
                hex_color = _rgb_to_hex(color)
                if hex_color:
                    return hex_color
            except Exception:
                pass
        try:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                solidFill = spPr.find(qn("a:solidFill"))
                if solidFill is not None:
                    hex_color = _extract_hex_from_color_choice(solidFill, theme_map, color_map)
                    if hex_color:
                        return hex_color
        except Exception:
            pass
    except Exception:
        pass
    return ""


def _extract_canvas_text_item_data(shape, slide, slide_w, slide_h, z_index, theme_map=None, color_map=None) -> dict | None:
    """Extract text from shape as a canvas text element with resolved font sizes and colors."""
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return None
    tf = shape.text_frame
    paragraphs = [p for p in tf.paragraphs if p.text.strip()]
    if not paragraphs:
        return None

    text_content = "\n".join(p.text for p in paragraphs)

    # Try to find a font size
    font_size = 16
    for p in paragraphs:
        for run in p.runs:
            if run.font.size:
                font_size = int(run.font.size / 12700)
                break
        if font_size != 16:
            break

    # Try to find a color
    color = "#111827"
    for p in paragraphs:
        for run in p.runs:
            hex_col = _extract_run_hex_color(run, p, shape, slide, theme_map, color_map)
            if hex_col:
                color = hex_col
                break
        if color != "#111827":
            break

    x = _emu_to_pct(shape.left, slide_w)
    y = _emu_to_pct(shape.top, slide_h)
    w = _emu_to_pct(shape.width, slide_w)
    h = _emu_to_pct(shape.height, slide_h)

    return {
        "id": _uid("canvasitem"),
        "type": "text",
        "zIndex": z_index,
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "rotation": 0,
        "color": color,
        "text": text_content,
        "fontSize": font_size,
        "fontFamily": "inherit",
        "fontWeight": "normal",
        "fontStyle": "normal",
        "textDecoration": "none",
        "textAlign": "left",
        "lineHeight": 1.5,
        "letterSpacing": 0,
        "boxBg": "#ffffff",
        "boxBgOpacity": 0,
    }


def _extract_table_as_plain_text(shape) -> str | None:
    """Extract a PPTX table as a clean pipe-separated plain text string."""
    try:
        table = shape.table
        lines = []
        for row in table.rows:
            cells = [cell.text_frame.text.strip() if cell.text_frame else "" for cell in row.cells]
            lines.append(" | ".join(cells))
        return "\n".join(lines)
    except Exception:
        return None


# ── Main entry point ─────────────────────────────────────────────────────────

async def extract_slides(file_bytes: bytes):
    prs = Presentation(io.BytesIO(file_bytes))

    # 1. Build the dynamic theme color map for this specific presentation
    default_theme_map = _build_theme_color_map(prs)

    # Slide dimensions in EMU — needed to convert shape positions to percentages.
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    blocks = []

    for i, slide in enumerate(prs.slides):
        slide_title = f"Slide {i + 1}"
        items = []
        z_index = 1
        color_map = _get_effective_color_map(slide)
        slide_theme_bundle = _get_cached_theme_bundle(_get_theme_part_for_slide(slide))
        theme_map = slide_theme_bundle.get("colors") or default_theme_map

        # ── Determine slide title from title placeholder ──────────────────────
        for shape in slide.shapes:
            if (
                getattr(shape, "is_placeholder", False)
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
                and hasattr(shape, "text")
                and shape.text.strip()
            ):
                slide_title = shape.text.strip()
                break

        # ── Background ────────────────────────────────────────────────────────
        # Use the slide's own master theme bundle so multi-master decks keep
        # their real accent/background colors instead of borrowing another theme.
        bg = _extract_slide_background(slide, theme_map, color_map, slide_theme_bundle)

        # ── Transition ────────────────────────────────────────────────────────
        transition = _extract_transition(slide)

        # ── Slide transition canvas item ──────────────────────────────────────
        if transition:
            items.append({
                "id": _uid("canvasitem"),
                "type": "text",
                "zIndex": 9999,
                "x": 2.0,
                "y": 2.0,
                "w": 30.0,
                "h": 6.0,
                "rotation": 0,
                "color": "#8b1a1a",
                "text": f"▶ TRANSITION: {transition}",
                "fontSize": 10,
                "fontFamily": "inherit",
                "fontWeight": "bold",
                "fontStyle": "normal",
                "textDecoration": "none",
                "textAlign": "left",
                "lineHeight": 1.5,
                "letterSpacing": 0,
                "boxBg": "#fff5f5",
                "boxBgOpacity": 100,
            })

        # ── Embedded audio ────────────────────────────────────────────────────
        audio_blocks = _extract_embedded_audio(slide, slide_w, slide_h)
        bg_audio = None
        if audio_blocks:
            first_audio = audio_blocks[0]
            bg_audio = {
                "url": first_audio["audioUrl"],
                "name": first_audio.get("label") or "Audio Track",
                "mediaId": first_audio.get("mediaId")
            }

        # ── Per-shape extraction ──────────────────────────────────────────────
        for shape in slide.shapes:
            # --- Table ---
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_text = _extract_table_as_plain_text(shape)
                if table_text:
                    items.append({
                        "id": _uid("canvasitem"),
                        "type": "text",
                        "zIndex": z_index,
                        "x": _emu_to_pct(shape.left, slide_w),
                        "y": _emu_to_pct(shape.top, slide_h),
                        "w": _emu_to_pct(shape.width, slide_w),
                        "h": _emu_to_pct(shape.height, slide_h),
                        "rotation": 0,
                        "color": "#111827",
                        "text": table_text,
                        "fontSize": 14,
                        "fontFamily": "inherit",
                        "fontWeight": "normal",
                        "fontStyle": "normal",
                        "textDecoration": "none",
                        "textAlign": "left",
                        "lineHeight": 1.5,
                        "letterSpacing": 0,
                        "boxBg": "#ffffff",
                        "boxBgOpacity": 0,
                    })
                    z_index += 1
                continue

            # --- Image / Picture ---
            image_blob = None
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if _find_audio_rid(shape._element):
                    continue
                image_blob = shape.image.blob
            elif (
                getattr(shape, "is_placeholder", False)
                and hasattr(shape, "image")
            ):
                try:
                    image_blob = shape.image.blob
                except (AttributeError, Exception):
                    pass

            if image_blob:
                compressed, content_type = _compress_image(image_blob)
                data_uri = _blob_to_data_uri(compressed, content_type)
                alt_text = ""
                try:
                    alt_text = shape.name or ""
                except Exception:
                    pass
                items.append({
                    "id": _uid("canvasitem"),
                    "type": "image",
                    "zIndex": z_index,
                    "x": _emu_to_pct(shape.left, slide_w),
                    "y": _emu_to_pct(shape.top, slide_h),
                    "w": _emu_to_pct(shape.width, slide_w),
                    "h": _emu_to_pct(shape.height, slide_h),
                    "rotation": 0,
                    "src": data_uri,
                })
                z_index += 1
                continue

            # --- Skip MEDIA shapes ---
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                continue

            # --- Shape / Text box ---
            is_geom_shape = False
            geom_shape_type = None
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    ast = shape.auto_shape_type
                    if ast in (MSO_SHAPE.RECTANGLE, MSO_SHAPE.ROUNDED_RECTANGLE):
                        geom_shape_type = 'rect'
                        is_geom_shape = True
                    elif ast == MSO_SHAPE.OVAL:
                        geom_shape_type = 'circle'
                        is_geom_shape = True
                    elif ast in (MSO_SHAPE.ISOSCELES_TRIANGLE, MSO_SHAPE.RIGHT_TRIANGLE):
                        geom_shape_type = 'triangle'
                        is_geom_shape = True
                    else:
                        # Fallback for other autoshapes
                        geom_shape_type = 'rect'
                        is_geom_shape = True
                elif shape.shape_type == MSO_SHAPE_TYPE.RECTANGLE:
                    geom_shape_type = 'rect'
                    is_geom_shape = True
            except Exception:
                pass

            has_text = False
            if hasattr(shape, "text_frame"):
                tf = shape.text_frame
                paragraphs = [p for p in tf.paragraphs if p.text.strip()]
                if paragraphs:
                    has_text = True

            if is_geom_shape:
                fill_color = _extract_shape_fill_color(shape, theme_map, color_map)
                if not fill_color:
                    fill_color = "#3b82f6"

                items.append({
                    "id": _uid("canvasitem"),
                    "type": geom_shape_type,
                    "zIndex": z_index,
                    "x": _emu_to_pct(shape.left, slide_w),
                    "y": _emu_to_pct(shape.top, slide_h),
                    "w": _emu_to_pct(shape.width, slide_w),
                    "h": _emu_to_pct(shape.height, slide_h),
                    "rotation": 0,
                    "color": fill_color,
                })

                if has_text:
                    text_item = _extract_canvas_text_item_data(shape, slide, slide_w, slide_h, z_index + 1, theme_map, color_map)
                    if text_item:
                        items.append(text_item)
                        z_index += 2
                    else:
                        z_index += 1
                else:
                    z_index += 1

            elif has_text:
                text_item = _extract_canvas_text_item_data(shape, slide, slide_w, slide_h, z_index, theme_map, color_map)
                if text_item:
                    items.append(text_item)
                    z_index += 1

        # ── Speaker notes → appended as a styled text box canvas item ────────
        notes_text = _extract_notes(slide)
        if notes_text:
            items.append({
                "id": _uid("canvasitem"),
                "type": "text",
                "zIndex": z_index,
                "x": 5.0,
                "y": 82.0,
                "w": 90.0,
                "h": 15.0,
                "rotation": 0,
                "color": "#555555",
                "text": f"SPEAKER NOTES:\n{notes_text}",
                "fontSize": 12,
                "fontFamily": "inherit",
                "fontWeight": "normal",
                "fontStyle": "normal",
                "textDecoration": "none",
                "textAlign": "left",
                "lineHeight": 1.5,
                "letterSpacing": 0,
                "boxBg": "#fff5f5",
                "boxBgOpacity": 100,
            })
            z_index += 1

        # ── Build slide block ─────────────────────────────────────────────────
        canvas_bg = "#ffffff"
        background = {"type": "color", "value": "#ffffff"}
        if bg["type"] in ("color", "image") and bg["value"]:
            background = bg
            if bg["type"] == "color":
                canvas_bg = bg["value"]
            elif bg["type"] == "image":
                canvas_bg = f"url('{bg['value']}')"

        slide_block = {
            "id": _uid("slide"),
            "type": "canvas",
            "title": slide_title,
            "slideNumber": i + 1,
            "canvasBg": canvas_bg,
            "background": background,
            "items": items,
        }

        if bg_audio:
            slide_block["bgAudio"] = bg_audio

        blocks.append(slide_block)

    return blocks
